from pathlib import Path

from docx import Document as DocxDocument
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from pypdf import PdfReader
from pptx import Presentation

from app.rag.embeddings import get_vector_store


def extract_text(file_path: str) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if suffix in {".doc", ".docx"}:
        doc = DocxDocument(file_path)
        return "\n".join(paragraph.text for paragraph in doc.paragraphs)

    if suffix in {".ppt", ".pptx"}:
        deck = Presentation(file_path)
        text_blocks = []
        for slide in deck.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_blocks.append(shape.text)
        return "\n".join(text_blocks)

    return Path(file_path).read_text(encoding="utf-8", errors="ignore")


def ingest_material(material_id: int, course_id: int, file_path: str) -> int:
    text = extract_text(file_path)
    splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=150)
    chunks = splitter.split_text(text)

    docs = [
        Document(
            page_content=chunk,
            metadata={"material_id": material_id, "course_id": course_id, "source": file_path, "chunk": index},
        )
        for index, chunk in enumerate(chunks)
        if chunk.strip()
    ]

    if not docs:
        return 0

    store = get_vector_store()
    store.add_documents(docs)
    return len(docs)
